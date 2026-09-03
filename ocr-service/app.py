"""OCR / large-document fallback service for the AKLA Matter Hub's document
pipeline.

extractText.ts (the shared module every document-processing edge function
uses) can only read a PDF's real text layer — scanned/photographed signed
copies have none, and come back essentially empty. Running real Tesseract
OCR isn't viable inside a Supabase Edge Function (150-256MB memory, 2s CPU
time per request — OCR blows through both), so this is a small dedicated
service on its own VM that extractText.ts calls as a fallback when normal
extraction yields ~nothing.

Poppler rasterizes the whole PDF to page image files on disk in one pass
(confirmed the hard way: doing this per-page instead re-parses the entire
source PDF from scratch on every single page — 54x redundant work for a
54-page document, enough to blow past a 2-minute timeout on its own before
Tesseract even starts). Tesseract then runs over those files one at a time,
deleting each as it goes — memory and disk both stay flat regardless of
document length.

/extract/docx exists for the same underlying reason, different failure
mode: mammoth's in-Edge-Function docx-to-HTML conversion needs to hold the
whole unzipped document.xml (plus mammoth's own intermediate structures) in
memory at once, and a large Word file with heavy tracked-changes history
(every insertion/deletion is its own XML run) can expand to many times its
zipped size — confirmed directly against a real 5.5MB precedent file, which
killed the Edge Function with Supabase's WORKER_RESOURCE_LIMIT in under 7
seconds (a hard memory kill, not a timeout — nothing in the function's own
code could have caught it). This VM has no such ceiling, so extractText.ts
routes docx files over a size threshold here instead of attempting mammoth
in-process at all.
"""

import io
import os
import tempfile

from flask import Flask, request, jsonify
from pdf2image import convert_from_path
from pdf2image.pdf2image import pdfinfo_from_path
from PIL import Image
import pytesseract
import mammoth
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = Flask(__name__)

OCR_SHARED_SECRET = os.environ["OCR_SHARED_SECRET"]
OCR_DPI = 150  # enough for typed contract text; keeps per-page size modest

# Cheap safety valve — an unbounded page count would just tie up the service
# for minutes on end; anything this long is unusual for a contract and worth
# a human's attention rather than a silent multi-minute hang.
MAX_PAGES = 200

# Same idea for docx — this VM has no per-request memory ceiling, but an
# absurdly large upload (a corrupt file, someone's entire drive zipped into
# a .docx) shouldn't be allowed to run unbounded either.
MAX_DOCX_BYTES = 50 * 1024 * 1024

# Decks are mostly images, so the bar is higher than docx — but still a bar.
MAX_PPTX_BYTES = 200 * 1024 * 1024


def _check_auth():
    auth = request.headers.get("Authorization", "")
    return auth == f"Bearer {OCR_SHARED_SECRET}"


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"ok": True})


@app.route("/ocr/pdf", methods=["POST"])
def ocr_pdf():
    if not _check_auth():
        return jsonify({"error": "Unauthorized"}), 401

    data = request.get_data()
    if not data:
        return jsonify({"error": "No file body provided"}), 400

    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = os.path.join(tmpdir, "input.pdf")
        with open(pdf_path, "wb") as f:
            f.write(data)

        try:
            page_count = pdfinfo_from_path(pdf_path)["Pages"]
        except Exception as err:
            return jsonify({"error": f"Could not read PDF: {err}"}), 400

        if page_count > MAX_PAGES:
            return jsonify({"error": f"PDF has {page_count} pages, over the {MAX_PAGES}-page OCR limit"}), 413

        image_paths = convert_from_path(
            pdf_path, dpi=OCR_DPI, output_folder=tmpdir, paths_only=True, fmt="ppm"
        )

        texts = []
        for image_path in sorted(image_paths):
            with Image.open(image_path) as img:
                texts.append(pytesseract.image_to_string(img))
            os.remove(image_path)

    return jsonify({"text": "\n\n".join(texts), "pages": page_count})


@app.route("/extract/docx", methods=["POST"])
def extract_docx():
    if not _check_auth():
        return jsonify({"error": "Unauthorized"}), 401

    data = request.get_data()
    if not data:
        return jsonify({"error": "No file body provided"}), 400

    if len(data) > MAX_DOCX_BYTES:
        return jsonify({"error": f"File is {len(data)} bytes, over the {MAX_DOCX_BYTES}-byte limit"}), 413

    try:
        result = mammoth.convert_to_html(io.BytesIO(data))
    except Exception as err:
        return jsonify({"error": f"Could not convert docx: {err}"}), 400

    return jsonify({"html": result.value})


def _text_frame_lines(text_frame, lines):
    for paragraph in text_frame.paragraphs:
        # paragraph.text renders an <a:br/> as a vertical tab.
        text = paragraph.text.replace("\v", "\n").strip()
        if text:
            lines.append(text)


def _shape_lines(shape, lines):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _shape_lines(child, lines)
        return
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        _text_frame_lines(shape.text_frame, lines)
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _text_frame_lines(cell.text_frame, lines)


@app.route("/extract/pptx", methods=["POST"])
def extract_pptx():
    """Large-deck counterpart of extractText.ts's in-process pptx path. The
    output shape is deliberately identical to that path — "[Slide N]" blocks,
    one line per paragraph in document order, notes appended as "Notes: …" —
    because redline suggestions quote the extracted text verbatim and the
    same file can go down either path depending only on its size."""
    if not _check_auth():
        return jsonify({"error": "Unauthorized"}), 401

    data = request.get_data()
    if not data:
        return jsonify({"error": "No file body provided"}), 400

    if len(data) > MAX_PPTX_BYTES:
        return jsonify({"error": f"File is {len(data)} bytes, over the {MAX_PPTX_BYTES}-byte limit"}), 413

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as err:
        return jsonify({"error": f"Could not open pptx: {err}"}), 400

    blocks = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = []
        for shape in slide.shapes:
            _shape_lines(shape, lines)
        if slide.has_notes_slide:
            note_lines = []
            _text_frame_lines(slide.notes_slide.notes_text_frame, note_lines)
            if note_lines:
                lines.append("Notes: " + "\n".join(note_lines))
        blocks.append(f"[Slide {index}]\n" + "\n".join(lines))

    return jsonify({"text": "\n\n".join(blocks), "slides": len(blocks)})


if __name__ == "__main__":
    app.run(host="127.0.0.1", port=8090)
